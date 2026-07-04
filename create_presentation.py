import os
import pptx
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# 1. Initialize presentation and set to Widescreen (16:9)
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6] # Blank layout

# 2. Design Theme & Colors
COLOR_DARK_SLATE = RGBColor(30, 41, 59)   # #1E293B
COLOR_LIGHT_GRAY = RGBColor(248, 250, 252) # #F8FAFC
COLOR_TEXT_DARK = RGBColor(15, 23, 42)     # #0F172A
COLOR_TEXT_LIGHT = RGBColor(255, 255, 255) # #FFFFFF
COLOR_TEAL = RGBColor(15, 118, 110)       # #0F766E
COLOR_CORAL = RGBColor(190, 18, 60)       # #BE123C
COLOR_MUTED = RGBColor(100, 116, 139)     # #64748B
COLOR_WHITE = RGBColor(255, 255, 255)

FONT_HEADING = "Segoe UI"
FONT_BODY = "Segoe UI"

# Paths to images
base_dir = r"C:\Users\visha\Desktop\ML\Projects\Customer_churn"
img_churn_dist = os.path.join(base_dir, "images", "churn_distribution.png")
img_pm_vs_churn = os.path.join(base_dir, "extracted_plots", "plot_cell_58_out_0.png")
img_int_vs_churn = os.path.join(base_dir, "extracted_plots", "plot_cell_63_out_0.png")
img_contract_vs_churn = os.path.join(base_dir, "extracted_plots", "plot_cell_73_out_0.png")
img_tenure_box = os.path.join(base_dir, "extracted_plots", "plot_cell_92_out_0.png")
img_charges_box = os.path.join(base_dir, "extracted_plots", "plot_cell_93_out_0.png")
img_heatmap = os.path.join(base_dir, "extracted_plots", "plot_cell_102_out_0.png")

def set_slide_background(slide, color):
    """Draws a full-slide rectangle to serve as a custom background color."""
    bg = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background() # No line
    # Send to back by moving it to the first position in the shape tree
    slide.shapes._spTree.remove(bg._element)
    slide.shapes._spTree.insert(2, bg._element)
    return bg

def add_header(slide, title_text):
    """Adds a standardized clean header zone for content slides."""
    # Header Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.7))
    tf = title_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.name = FONT_HEADING
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEXT_DARK

    # Teal accent line under title
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2.0), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TEAL
    line.line.fill.background()

def create_title_slide():
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_DARK_SLATE)
    
    # Left Teal Accent Bar
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(2.2), Inches(0.15), Inches(3.0)
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_TEAL
    bar.line.fill.background()
    
    # Text Block
    text_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.1), Inches(11.0), Inches(3.2))
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    # Main Title
    p1 = tf.paragraphs[0]
    p1.text = "Customer Churn Analysis"
    p1.font.name = FONT_HEADING
    p1.font.size = Pt(46)
    p1.font.bold = True
    p1.font.color.rgb = COLOR_TEXT_LIGHT
    
    p1_sub = tf.add_paragraph()
    p1_sub.text = "& Machine Learning Predictions"
    p1_sub.font.name = FONT_HEADING
    p1_sub.font.size = Pt(36)
    p1_sub.font.bold = True
    p1_sub.font.color.rgb = COLOR_TEAL
    p1_sub.space_after = Pt(24)
    
    # Subtitle
    p2 = tf.add_paragraph()
    p2.text = "Predicting customer retention, identifying key churn drivers, and proposing data-driven strategies."
    p2.font.name = FONT_BODY
    p2.font.size = Pt(16)
    p2.font.color.rgb = COLOR_MUTED
    
    p3 = tf.add_paragraph()
    p3.text = "Dataset: Telco Customer Churn (7,043 Accounts)"
    p3.font.name = FONT_BODY
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_TEAL
    p3.space_before = Pt(30)

def create_slide_2():
    """Executive Summary & Objective"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_LIGHT_GRAY)
    add_header(slide, "Executive Summary & Project Goal")
    
    # Left Text Column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_head = tf.paragraphs[0]
    p_head.text = "The Business Problem"
    p_head.font.name = FONT_HEADING
    p_head.font.size = Pt(20)
    p_head.font.bold = True
    p_head.font.color.rgb = COLOR_TEAL
    p_head.space_after = Pt(12)
    
    bullets = [
        ("Customer Churn", " occurs when subscribers cancel their service. Acquiring new customers costs 5-25x more than retaining existing ones."),
        ("Key Objectives", ": Analyze customer behavior, clean and preprocess customer data, train predictive models (Logistic Regression, Decision Trees, Random Forests), and determine the primary risk factors."),
        ("Base Dataset", ": 7,043 unique customer accounts with 21 socio-demographic, service subscription, and financial attributes."),
        ("Base Churn Rate", ": 26.54% of customers have churned (1,869 churned vs. 5,174 active). This represents a significant revenue leakage.")
    ]
    
    for title, desc in bullets:
        p = tf.add_paragraph()
        p.space_after = Pt(14)
        run_title = p.add_run()
        run_title.text = "• " + title
        run_title.font.bold = True
        run_title.font.size = Pt(15)
        run_title.font.color.rgb = COLOR_TEXT_DARK
        
        run_desc = p.add_run()
        run_desc.text = desc
        run_desc.font.size = Pt(15)
        run_desc.font.color.rgb = COLOR_MUTED
        if "26.54%" in desc:
            # Highlight this stat
            run_desc.font.color.rgb = COLOR_CORAL
            run_desc.font.bold = True

    # Right Image
    if os.path.exists(img_churn_dist):
        slide.shapes.add_picture(img_churn_dist, Inches(7.2), Inches(1.8), width=Inches(5.3), height=Inches(4.5))

def create_slide_3():
    """Data Preprocessing & Cleaning"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_LIGHT_GRAY)
    add_header(slide, "Data Pipeline: Preprocessing & Engineering")
    
    # Left Column (Cleaning)
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    pl = tf_l.paragraphs[0]
    pl.text = "1. Data Cleaning & Structuring"
    pl.font.name = FONT_HEADING
    pl.font.size = Pt(20)
    pl.font.bold = True
    pl.font.color.rgb = COLOR_TEAL
    pl.space_after = Pt(12)
    
    cleaning_steps = [
        ("ID Removal", "Dropped customerID as it has no predictive power and adds noise."),
        ("Type Conversion", "TotalCharges contained blank space strings. These were coerced to numeric (NaNs generated)."),
        ("Imputation", "Filled missing TotalCharges values (11 records) using the column mean value."),
        ("Deduplication", "Identified and removed duplicate customer records to ensure model integrity.")
    ]
    
    for title, desc in cleaning_steps:
        p = tf_l.add_paragraph()
        p.space_after = Pt(12)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "• " + title + ": "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT_DARK
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_MUTED

    # Right Column (Feature Engineering)
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.2))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0
    
    pr = tf_r.paragraphs[0]
    pr.text = "2. Feature Engineering & Scaling"
    pr.font.name = FONT_HEADING
    pr.font.size = Pt(20)
    pr.font.bold = True
    pr.font.color.rgb = COLOR_TEAL
    pr.space_after = Pt(12)
    
    eng_steps = [
        ("Binary Mapping", "Mapped binary categorical variables (gender, Partner, Dependents, PhoneService, PaperlessBilling, Churn) to 0 and 1 (e.g., Male=1, Female=0)."),
        ("Dummy Encoding", "Applied pd.get_dummies(drop_first=True) to multi-class variables (Contract, PaymentMethod, InternetService, and protection add-ons)."),
        ("Train/Test Split", "Divided dataset into 80% training set (to fit models) and 20% test set (to evaluate generalization)."),
        ("Feature Scaling", "Standardized features using StandardScaler on training data and transformed test data to prevent scale bias.")
    ]
    
    for title, desc in eng_steps:
        p = tf_r.add_paragraph()
        p.space_after = Pt(12)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "• " + title + ": "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT_DARK
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_MUTED

def create_slide_4():
    """EDA - Demographics & Financials"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_LIGHT_GRAY)
    add_header(slide, "EDA: Customer Tenure & Monthly Charges")
    
    # Left Text Column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_head = tf.paragraphs[0]
    p_head.text = "Distribution Insights (Boxplots)"
    p_head.font.name = FONT_HEADING
    p_head.font.size = Pt(20)
    p_head.font.bold = True
    p_head.font.color.rgb = COLOR_TEAL
    p_head.space_after = Pt(12)
    
    insights = [
        ("Tenure Distribution", " (Top Chart): Shows a bimodal distribution. A large group of customers have very short tenure (0-10 months) and represent the highest churn risk, while another large group is highly stable (60-72 months)."),
        ("Monthly Charges", " (Bottom Chart): The boxplot highlights that monthly charges are concentrated between $35 and $90. Higher monthly bills are directly linked with a higher churn probability."),
        ("Key Takeaway", ": The risk of churn is heavily concentrated in the early months of a customer's life cycle, particularly if they are billed high monthly rates immediately.")
    ]
    
    for title, desc in insights:
        p = tf.add_paragraph()
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "• " + title + ": "
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_TEXT_DARK
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = COLOR_MUTED
        if "early months" in desc:
            r2.font.color.rgb = COLOR_CORAL
            r2.font.bold = True

    # Right Images
    if os.path.exists(img_tenure_box):
        slide.shapes.add_picture(img_tenure_box, Inches(7.2), Inches(1.5), width=Inches(5.3), height=Inches(2.5))
    if os.path.exists(img_charges_box):
        slide.shapes.add_picture(img_charges_box, Inches(7.2), Inches(4.3), width=Inches(5.3), height=Inches(2.5))

def create_slide_5():
    """EDA - Services & Contract Types"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_LIGHT_GRAY)
    add_header(slide, "EDA: Services & Contract Types vs. Churn")
    
    # Left Text Column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_head = tf.paragraphs[0]
    p_head.text = "Service & Contract Risk Factors"
    p_head.font.name = FONT_HEADING
    p_head.font.size = Pt(20)
    p_head.font.bold = True
    p_head.font.color.rgb = COLOR_TEAL
    p_head.space_after = Pt(12)
    
    insights = [
        ("Contract Type Influence", ": Month-to-month contracts have an extremely high churn rate. Conversely, customers signed on 1-year and 2-year contracts show remarkable stability and very low churn rates."),
        ("Internet Service Risk", ": Fiber Optic subscribers exhibit a significantly higher rate of churn compared to DSL or No-Internet users. This indicates issues with price sensitivity, service quality, or heavy competitor targeting."),
        ("Payment Methods", ": Electronic Check is highly correlated with churn. Auto-pay payment methods (Credit Card, Bank Transfer) display high retention, representing a major optimization area.")
    ]
    
    for title, desc in insights:
        p = tf.add_paragraph()
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "• " + title
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_TEXT_DARK
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = COLOR_MUTED
        if "Month-to-month" in desc or "Fiber Optic" in desc:
            r2.font.bold = True
            r2.font.color.rgb = COLOR_CORAL

    # Right Images
    if os.path.exists(img_contract_vs_churn):
        slide.shapes.add_picture(img_contract_vs_churn, Inches(7.2), Inches(1.5), width=Inches(5.3), height=Inches(2.5))
    if os.path.exists(img_int_vs_churn):
        slide.shapes.add_picture(img_int_vs_churn, Inches(7.2), Inches(4.3), width=Inches(5.3), height=Inches(2.5))

def create_slide_6():
    """Correlation Heatmap"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_LIGHT_GRAY)
    add_header(slide, "Correlation Heatmap & Feature Inter-relationships")
    
    # Left Text Column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.0), Inches(5.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_head = tf.paragraphs[0]
    p_head.text = "Heatmap Analysis"
    p_head.font.name = FONT_HEADING
    p_head.font.size = Pt(20)
    p_head.font.bold = True
    p_head.font.color.rgb = COLOR_TEAL
    p_head.space_after = Pt(12)
    
    insights = [
        ("Internet Services Link", "The heatmap reveals moderate-to-strong relationships among internet-related service features, indicating that customers often subscribe to multiple value-added services together."),
        ("Multicollinearity", "Detected between tenure, MonthlyCharges, and TotalCharges (since TotalCharges is a mathematical product of the other two). This must be handled during modeling (e.g. through regularized algorithms or scale adjustments)."),
        ("Strongest Direct Churn Correlations", "Internet Service (Fiber Optic) and Payment Method (Electronic Check) show the highest positive correlation coefficients with the target variable Churn.")
    ]
    
    for title, desc in insights:
        p = tf.add_paragraph()
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "• " + title + ": "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEXT_DARK
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_MUTED

    # Right Image
    if os.path.exists(img_heatmap):
        slide.shapes.add_picture(img_heatmap, Inches(6.2), Inches(1.5), width=Inches(6.3), height=Inches(5.0))

def create_slide_7():
    """Model Selection & Training"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_LIGHT_GRAY)
    add_header(slide, "Model Development & Validation Strategy")
    
    # Left Text Column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(5.2))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p_head = tf.paragraphs[0]
    p_head.text = "Evaluation Framework & Algorithms"
    p_head.font.name = FONT_HEADING
    p_head.font.size = Pt(20)
    p_head.font.bold = True
    p_head.font.color.rgb = COLOR_TEAL
    p_head.space_after = Pt(12)
    
    models = [
        ("Logistic Regression", "Used as our baseline model. Standardized inputs, evaluated with cross-validation. Shows stable convergence, high interpretability, and excellent ROC-AUC performance."),
        ("Decision Tree Classifier", "Initially fit without constraints (resulting in heavy overfitting, 1.0 training accuracy). Subsequently regularized with max_depth=5 to limit splits, improving generalization on unseen test data."),
        ("Random Forest Classifier", "An ensemble method trained with 100 estimators. Bagging reduces variance and provides robust predictions, though it suffers slightly in test accuracy here compared to Logistic Regression."),
        ("Validation Protocol", "Divided dataset into 80% train / 20% test splits. Features normalized using StandardScaler. Evaluated Logistic Regression stability using 5-Fold Cross Validation.")
    ]
    
    for name, desc in models:
        p = tf.add_paragraph()
        p.space_after = Pt(14)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "• " + name + ": "
        r1.font.bold = True
        r1.font.size = Pt(15)
        r1.font.color.rgb = COLOR_TEXT_DARK
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = COLOR_MUTED

def create_slide_8():
    """Model Evaluation & Results Table"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_LIGHT_GRAY)
    add_header(slide, "Model Performance Comparison")
    
    # Left Text Column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(11.733), Inches(1.8))
    tf = left_box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    
    p = tf.paragraphs[0]
    p.text = "Results Summary"
    p.font.name = FONT_HEADING
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = COLOR_TEAL
    p.space_after = Pt(6)
    
    p_desc = tf.add_paragraph()
    p_desc.text = "Logistic Regression outperformed the tree-based models across overall metrics. While Decision Trees and Random Forests had slightly higher precision for churn prediction, Logistic Regression captured significantly more actual churners (higher recall) and showed superior overall discrimination (ROC-AUC: 0.862)."
    p_desc.font.name = FONT_BODY
    p_desc.font.size = Pt(14)
    p_desc.font.color.rgb = COLOR_MUTED
    p_desc.line_spacing = 1.15
    
    # Add native table
    rows = 4
    cols = 6
    left = Inches(0.8)
    top = Inches(3.2)
    width = Inches(11.733)
    height = Inches(2.2)
    
    table_shape = slide.shapes.add_table(rows, cols, left, top, width, height)
    table = table_shape.table
    
    # Table headers
    headers = ["Model", "Test Accuracy", "Churn Precision", "Churn Recall", "Churn F1-Score", "ROC-AUC"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_TEAL
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_HEADING
            p.font.size = Pt(14)
            p.font.bold = True
            p.font.color.rgb = COLOR_WHITE
            p.alignment = PP_ALIGN.CENTER
            
    # Data rows
    # 0 = Churn (1)
    data = [
        ["Logistic Regression", "82.26%", "0.69", "0.60", "0.64", "0.8620"],
        ["Decision Tree (max_depth=5)", "80.62%", "0.70", "0.46", "0.56", "-"],
        ["Random Forest (n=100)", "79.42%", "0.65", "0.46", "0.54", "-"]
    ]
    
    for row_idx, row_data in enumerate(data):
        for col_idx, val in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = val
            cell.fill.solid()
            # Highlight best model
            if row_idx == 0:
                cell.fill.fore_color.rgb = RGBColor(224, 242, 254) # Light blue accent
            else:
                cell.fill.fore_color.rgb = COLOR_WHITE
                
            for p in cell.text_frame.paragraphs:
                p.font.name = FONT_BODY
                p.font.size = Pt(13)
                p.alignment = PP_ALIGN.CENTER
                if row_idx == 0:
                    p.font.bold = True
                    p.font.color.rgb = COLOR_TEAL
                else:
                    p.font.color.rgb = COLOR_TEXT_DARK

def create_slide_9():
    """Key Drivers of Churn"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_LIGHT_GRAY)
    add_header(slide, "Key Drivers of Churn: Model Coefficients")
    
    # Left Column (Positive Drivers - Risk Factors)
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(5.6), Inches(5.2))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    pl = tf_l.paragraphs[0]
    pl.text = "Top Churn Drivers (Increases Risk)"
    pl.font.name = FONT_HEADING
    pl.font.size = Pt(20)
    pl.font.bold = True
    pl.font.color.rgb = COLOR_CORAL
    pl.space_after = Pt(12)
    
    pos_drivers = [
        ("Fiber Optic Internet (+0.625)", "The strongest positive predictor. Indicates high billing dissatisfaction, technical stability issues, or heavy competitor marketing."),
        ("Total Charges (+0.609)", "High absolute accumulated charges increase churn probability, likely capturing pricing pain points for older cohorts."),
        ("Streaming Movies / TV (+0.231 / +0.181)", "Premium entertainment services marginally increase churn likelihood, possibly due to cord-cutting behavior."),
        ("Paperless Billing (+0.164)", "Direct correlation, potentially due to lower friction in billing reviews or digital payments."),
        ("Electronic Check Payment (+0.151)", "Electronic check users churn significantly more than auto-pay users.")
    ]
    
    for title, desc in pos_drivers:
        p = tf_l.add_paragraph()
        p.space_after = Pt(10)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "▲ " + title + ": "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_CORAL
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_MUTED

    # Right Column (Negative Drivers - Loyalty Factors)
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.5), Inches(5.6), Inches(5.2))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0
    
    pr = tf_r.paragraphs[0]
    pr.text = "Top Loyalty Drivers (Reduces Risk)"
    pr.font.name = FONT_HEADING
    pr.font.size = Pt(20)
    pr.font.bold = True
    pr.font.color.rgb = COLOR_TEAL
    pr.space_after = Pt(12)
    
    neg_drivers = [
        ("Tenure Length (-1.310)", "The single strongest predictor of retention. The longer a customer stays, the dramatically less likely they are to leave."),
        ("Monthly Charges (-0.629)", "Negative coefficient (when accounting for tenure/total charges), reflecting high customer value lock-in."),
        ("Two-Year Contract (-0.619)", "Long-term contract commitments are highly effective in ensuring customer loyalty and preventing churn."),
        ("One-Year Contract (-0.269)", "Also provides moderate retention security, though less than two-year contracts."),
        ("Online Security / Tech Support (-0.158 / -0.120)", "Adding cybersecurity and tech assistance heavily locks customers in and reduces churn.")
    ]
    
    for title, desc in neg_drivers:
        p = tf_r.add_paragraph()
        p.space_after = Pt(10)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "▼ " + title + ": "
        r1.font.bold = True
        r1.font.size = Pt(14)
        r1.font.color.rgb = COLOR_TEAL
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(14)
        r2.font.color.rgb = COLOR_MUTED

def create_slide_10():
    """Actionable Recommendations (Dark Theme)"""
    slide = prs.slides.add_slide(blank_layout)
    set_slide_background(slide, COLOR_DARK_SLATE)
    
    # Header Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.733), Inches(0.7))
    tf_title = title_box.text_frame
    tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
    p_t = tf_title.paragraphs[0]
    p_t.text = "Strategic Recommendations for Business Impact"
    p_t.font.name = FONT_HEADING
    p_t.font.size = Pt(28)
    p_t.font.bold = True
    p_t.font.color.rgb = COLOR_TEXT_LIGHT

    # Teal accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.1), Inches(2.0), Inches(0.04)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = COLOR_TEAL
    line.line.fill.background()
    
    # Left Column
    left_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_l = left_box.text_frame
    tf_l.word_wrap = True
    tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = tf_l.margin_bottom = 0
    
    recs_l = [
        ("Migrate to Auto-Pay", "Provide a $2-$5 monthly discount for transitioning from Electronic Check to Credit Card or Bank Auto-Pay to automate renewals and reduce transaction friction."),
        ("Optimize Fiber Optic Retention", "Run quality audits and price-sensitivity check campaigns on Fiber Optic users. Bundle streaming discounts to cushion the premium price points.")
    ]
    
    for title, desc in recs_l:
        p = tf_l.add_paragraph()
        p.space_after = Pt(20)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "✔ " + title + "\n"
        r1.font.bold = True
        r1.font.size = Pt(18)
        r1.font.color.rgb = COLOR_TEAL
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = COLOR_MUTED

    # Right Column
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.6), Inches(5.6), Inches(5.2))
    tf_r = right_box.text_frame
    tf_r.word_wrap = True
    tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = tf_r.margin_bottom = 0
    
    recs_r = [
        ("Contract Conversion Incentives", "Target month-to-month subscribers with custom offers to upgrade to 1 or 2-year contracts. High risk of churn occurs in the first 12 months; lock-ins protect this window."),
        ("Bundle Security & Technical Add-ons", "Promote Online Security and Tech Support as a combined value pack. These sticky services increase retention and boost long-term customer life time value (LTV).")
    ]
    
    for title, desc in recs_r:
        p = tf_r.add_paragraph()
        p.space_after = Pt(20)
        p.line_spacing = 1.15
        
        r1 = p.add_run()
        r1.text = "✔ " + title + "\n"
        r1.font.bold = True
        r1.font.size = Pt(18)
        r1.font.color.rgb = COLOR_TEAL
        
        r2 = p.add_run()
        r2.text = desc
        r2.font.size = Pt(15)
        r2.font.color.rgb = COLOR_MUTED

# Generate the slides
create_title_slide()
create_slide_2()
create_slide_3()
create_slide_4()
create_slide_5()
create_slide_6()
create_slide_7()
create_slide_8()
create_slide_9()
create_slide_10()

# Save presentation
out_pptx_path = os.path.join(base_dir, "customer_churn_analysis.pptx")
prs.save(out_pptx_path)
print(f"Presentation saved successfully to: {out_pptx_path}")
